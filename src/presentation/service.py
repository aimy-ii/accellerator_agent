"""Презентация проекта.

Бизнес-логика заложена целиком, включается флагом PRESENTATION_PROVIDER:
    off   — узел пропускается (дефолт: сначала обкатываем ТЗ и подбор);
    local — LLM отдаёт структуру слайдов, .pptx собирается локально;
    gamma — структура уходит в Gamma API, оттуда приходит готовый .pptx.

Важно про Gamma: export_url (ссылка на .pptx) — ВРЕМЕННАЯ. Файл нужно
скачать сразу, пока ссылка жива, и положить к себе. presentation_url
(страница Gamma) — постоянная, храним её как preview_url.
"""
from __future__ import annotations

import asyncio
import io
import logging
import time

import httpx
from langchain_core.messages import HumanMessage, SystemMessage

from src.api.client import AcceleratorAPI
from src.core.config import settings
from src.presentation.models import Deck, PresentationResult
from src.presentation.prompts import DECK_SYSTEM, deck_user_message
from src.utils.llm_gen import ainvoke_llm, get_llm

log = logging.getLogger(__name__)


async def generate_deck(spec_text: str, roles: list[str]) -> Deck:
    """Генерирует структуру презентации по ТЗ."""
    async with get_llm(temperature=0.3) as llm:
        structured = llm.with_structured_output(Deck)
        deck: Deck = await ainvoke_llm(
            structured,
            [
                SystemMessage(content=DECK_SYSTEM),
                HumanMessage(content=deck_user_message(spec_text, roles)),
            ],
        )
    log.info("Структура презентации: слайдов=%d", len(deck.slides))
    return deck


# ─── Gamma ──────────────────────────────────────────────────────────────────

def _deck_to_markdown(deck: Deck) -> str:
    """Плоский markdown для Gamma: слайды разделяются '---'."""
    blocks = [f"# {deck.title}"]
    if deck.subtitle:
        blocks[0] += f"\n\n{deck.subtitle}"
    for slide in deck.slides:
        lines = [f"## {slide.title}"]
        lines.extend(f"- {b}" for b in slide.bullets)
        blocks.append("\n".join(lines))
    return "\n\n---\n\n".join(blocks)


async def _gamma_generate(deck: Deck) -> tuple[str | None, str | None]:
    """Гоняет структуру через Gamma API.

    Returns:
        (export_url, presentation_url). export_url — ВРЕМЕННАЯ ссылка на .pptx.
    """
    if not settings.gamma_api_key:
        raise RuntimeError("PRESENTATION_PROVIDER=gamma, но GAMMA_API_KEY не задан")

    headers = {
        "X-API-KEY": settings.gamma_api_key,
        "Content-Type": "application/json",
    }
    payload = {
        "inputText": _deck_to_markdown(deck),
        "textMode": "preserve",
        "format": "presentation",
        "exportAs": settings.gamma_export_as,
    }

    async with httpx.AsyncClient(timeout=60.0, trust_env=False) as client:
        created = await client.post(
            f"{settings.gamma_base_url}/generations",
            json=payload,
            headers=headers,
        )
        created.raise_for_status()
        generation_id = created.json().get("generationId")
        if not generation_id:
            raise RuntimeError(f"Gamma не вернула generationId: {created.text[:300]}")

        deadline = time.monotonic() + settings.gamma_poll_timeout
        while time.monotonic() < deadline:
            await asyncio.sleep(settings.gamma_poll_interval)
            status = await client.get(
                f"{settings.gamma_base_url}/generations/{generation_id}",
                headers=headers,
            )
            status.raise_for_status()
            data = status.json()
            state = (data.get("status") or "").lower()
            if state in ("completed", "success", "done"):
                return data.get("exportUrl"), data.get("gammaUrl")
            if state in ("failed", "error"):
                raise RuntimeError(f"Gamma вернула ошибку: {data}")

    raise TimeoutError("Gamma не успела сгенерировать презентацию за отведённое время")


async def _download(url: str) -> bytes:
    """Скачивает файл по временной ссылке Gamma."""
    async with httpx.AsyncClient(timeout=120.0, trust_env=False, follow_redirects=True) as client:
        response = await client.get(url)
        response.raise_for_status()
        return response.content


# ─── local (python-pptx) ────────────────────────────────────────────────────

def _render_pptx(deck: Deck) -> bytes:
    """Собирает .pptx локально из структуры слайдов."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    first = prs.slides.add_slide(title_layout)
    first.shapes.title.text = deck.title
    if deck.subtitle and len(first.placeholders) > 1:
        first.placeholders[1].text = deck.subtitle

    bullet_layout = prs.slide_layouts[1]
    for slide in deck.slides:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = slide.title
        body = s.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(slide.bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = bullet
            para.font.size = Pt(18)
        if slide.notes:
            s.notes_slide.notes_text_frame.text = slide.notes

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ─── точка входа ────────────────────────────────────────────────────────────

async def build_presentation(
    spec_text: str,
    roles: list[str],
    project_id: int,
    api: AcceleratorAPI,
    *,
    file_stem: str = "presentation",
) -> PresentationResult:
    """Генерирует презентацию и крепит .pptx к проекту.

    Провайдер берётся из settings.presentation_provider.
    При off — возвращает пустой результат, узел графа просто пропускается.
    """
    provider = (settings.presentation_provider or "off").lower()
    if provider == "off":
        return PresentationResult(provider="off")

    deck = await generate_deck(spec_text, roles)
    file_name = f"{file_stem}.pptx"

    try:
        if provider == "gamma":
            export_url, preview_url = await _gamma_generate(deck)
            if not export_url:
                return PresentationResult(
                    provider="gamma",
                    deck=deck,
                    preview_url=preview_url,
                    error="Gamma не вернула ссылку на файл",
                )
            # Ссылка временная — качаем немедленно, пока жива.
            content = await _download(export_url)
        else:
            content = _render_pptx(deck)
            preview_url = None

        await api.attach_file_to_project(
            project_id,
            file_name,
            content,
            mime_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
        )

        project = await api.get_project(project_id)
        file_url = next(
            (
                f["file_url"]
                for f in (project or {}).get("files", [])
                if f.get("file_name") == file_name
            ),
            None,
        )

        return PresentationResult(
            provider=provider,
            deck=deck,
            file_url=file_url,
            preview_url=preview_url,
        )
    except Exception as exc:  # noqa: BLE001
        log.error("Презентация не собралась: %s", exc)
        return PresentationResult(provider=provider, deck=deck, error=str(exc))
